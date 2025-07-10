#!/usr/bin/env python3
import argparse
import json
import os
import subprocess
import sys
import tempfile
from collections import defaultdict

STATUS_ORDER = [
    'Added',
    'Removed',
    'Renamed',
    'Modified',
    'RenamedAndModified',
    'Unchanged',
]

ATTR_CONTENT = """*.docx diff=custom
*.doc diff=custom
*.xlsx diff=custom
*.xls diff=custom
*.pptx diff=custom
*.ppt diff=custom
*.pdf diff=custom
"""


def conv_docx(path):
    import docx2txt
    return docx2txt.process(path)


def conv_doc(path):
    out = subprocess.run(['strings', path], stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    return out.stdout


def conv_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    text = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            text.append('\t'.join([str(c) if c is not None else '' for c in row]))
    return '\n'.join(text)


def conv_xls(path):
    import xlrd
    wb = xlrd.open_workbook(path)
    text = []
    for sheet in wb.sheets():
        for rx in range(sheet.nrows):
            text.append('\t'.join([str(cell.value) for cell in sheet.row(rx)]))
    return '\n'.join(text)


def conv_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text.append(shape.text)
    return '\n'.join(text)


def conv_ppt(path):
    out = subprocess.run(['strings', path], stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    return out.stdout


def conv_pdf(path):
    from PyPDF2 import PdfReader
    reader = PdfReader(path)
    text = []
    for page in reader.pages:
        text.append(page.extract_text() or '')
    return '\n'.join(text)


_MAPPING = {
    '.docx': conv_docx,
    '.doc': conv_doc,
    '.xlsx': conv_xlsx,
    '.xls': conv_xls,
    '.pptx': conv_pptx,
    '.ppt': conv_ppt,
    '.pdf': conv_pdf,
}


def textconv_main(path: str) -> None:
    ext = os.path.splitext(path)[1].lower()
    fn = _MAPPING.get(ext)
    if fn:
        try:
            text = fn(path)
            sys.stdout.write(text)
        except Exception:
            pass
    else:
        with open(path, 'r', errors='ignore') as f:
            sys.stdout.write(f.read())


def run_git_diff(dir1, dir2, attrs, textconv_path):
    with tempfile.NamedTemporaryFile('w', delete=False) as f:
        f.write(attrs)
        attr_file = f.name
    abs1 = os.path.abspath(dir1)
    abs2 = os.path.abspath(dir2)
    cmd = [
        'git',
        '-c', f'diff.custom.textconv={textconv_path}',
        '-c', f'core.attributesfile={attr_file}',
        'diff', '--no-index', '--textconv', '-M', '--name-status', abs1, abs2
    ]
    result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    os.unlink(attr_file)
    return result.stdout


def get_diff_hunk(path1, path2, attrs, textconv_path):
    with tempfile.NamedTemporaryFile('w', delete=False) as f:
        f.write(attrs)
        attr_file = f.name
    abs1 = os.path.abspath(path1)
    abs2 = os.path.abspath(path2)
    cmd = [
        'git',
        '-c', f'diff.custom.textconv={textconv_path}',
        '-c', f'core.attributesfile={attr_file}',
        'diff', '--no-index', '--textconv', '-M', abs1, abs2
    ]
    result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    os.unlink(attr_file)
    return result.stdout


def collect_files(root):
    files = []
    for base, _, names in os.walk(root):
        for n in names:
            files.append(os.path.relpath(os.path.join(base, n), root))
    return set(files)


def main():
    parser = argparse.ArgumentParser(description='Compare directories using git diff with textconv.')
    parser.add_argument('dir1')
    parser.add_argument('dir2')
    parser.add_argument('--json', action='store_true', help='Output in JSON format')
    args = parser.parse_args()

    dir1 = os.path.abspath(args.dir1)
    dir2 = os.path.abspath(args.dir2)

    diff_output = run_git_diff(dir1, dir2, ATTR_CONTENT, os.path.abspath(__file__))
    changes = []
    changed_paths = set()

    for line in diff_output.strip().splitlines():
        if not line:
            continue
        parts = line.split('\t')
        status = parts[0]
        if status.startswith('R'):
            score = int(status[1:])
            old = os.path.relpath(parts[1], dir1)
            new = os.path.relpath(parts[2], dir2)
            if score == 100:
                st = 'Renamed'
            else:
                st = 'RenamedAndModified'
            changes.append({'status': st, 'old': old, 'new': new})
            changed_paths.add(old)
            changed_paths.add(new)
        elif status == 'M':
            path = os.path.relpath(parts[1], dir1)
            changes.append({'status': 'Modified', 'path': path})
            changed_paths.add(path)
        elif status == 'A':
            path = os.path.relpath(parts[1], dir2)
            changes.append({'status': 'Added', 'path': path})
            changed_paths.add(path)
        elif status == 'D':
            path = os.path.relpath(parts[1], dir1)
            changes.append({'status': 'Removed', 'path': path})
            changed_paths.add(path)

    all_files_dir1 = collect_files(dir1)
    all_files_dir2 = collect_files(dir2)
    unchanged = sorted((all_files_dir1 & all_files_dir2) - changed_paths)
    for p in unchanged:
        changes.append({'status': 'Unchanged', 'path': p})

    # get diffs for changed files
    for change in changes:
        st = change['status']
        if st in ('Modified', 'RenamedAndModified'):
            old = os.path.join(dir1, change.get('old', change.get('path')))
            new = os.path.join(dir2, change.get('new', change.get('path')))
            change['diff'] = get_diff_hunk(old, new, ATTR_CONTENT, os.path.abspath(__file__))
        elif st == 'Added':
            new = os.path.join(dir2, change['path'])
            change['diff'] = get_diff_hunk(os.devnull, new, ATTR_CONTENT, os.path.abspath(__file__))
        elif st == 'Removed':
            old = os.path.join(dir1, change['path'])
            change['diff'] = get_diff_hunk(old, os.devnull, ATTR_CONTENT, os.path.abspath(__file__))
        elif st == 'Renamed':
            old = os.path.join(dir1, change['old'])
            new = os.path.join(dir2, change['new'])
            change['diff'] = get_diff_hunk(old, new, ATTR_CONTENT, os.path.abspath(__file__))

    # group by status order
    grouped = defaultdict(list)
    for c in changes:
        grouped[c['status']].append(c)

    if args.json:
        out = {st: grouped.get(st, []) for st in STATUS_ORDER}
        print(json.dumps(out, indent=2, ensure_ascii=False))
    else:
        for st in STATUS_ORDER:
            items = grouped.get(st)
            if not items:
                continue
            print(f'## {st}')
            for item in items:
                if st in ('Renamed', 'RenamedAndModified'):
                    print(f"{item['old']} -> {item['new']}")
                else:
                    print(item.get('path'))
                if 'diff' in item:
                    print(item['diff'])

if __name__ == '__main__':
    if len(sys.argv) == 2 and os.path.isfile(sys.argv[1]) and not os.path.isdir(sys.argv[1]):
        textconv_main(sys.argv[1])
    else:
        main()

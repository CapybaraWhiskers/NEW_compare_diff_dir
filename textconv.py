#!/usr/bin/env python3
import sys, os, subprocess

def conv_docx(path):
    import docx2txt
    return docx2txt.process(path)

def conv_doc(path):
    # Use strings as fallback
    out = subprocess.run(['strings', path], stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    return out.stdout

def conv_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    text = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            text.append('\t'.join([str(c) if c is not None else '' for c in row]))
    return '\n'.join(text)

def conv_xls(path):
    import xlrd
    wb = xlrd.open_workbook(path)
    text = []
    for sheet in wb.sheets():
        for rx in range(sheet.nrows):
            text.append('\t'.join([str(cell.value) for cell in sheet.row(rx)]))
    return '\n'.join(text)

def conv_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text.append(shape.text)
    return '\n'.join(text)

def conv_ppt(path):
    out = subprocess.run(['strings', path], stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    return out.stdout

def conv_pdf(path):
    from PyPDF2 import PdfReader
    reader = PdfReader(path)
    text = []
    for page in reader.pages:
        text.append(page.extract_text() or '')
    return '\n'.join(text)

mapping = {
    '.docx': conv_docx,
    '.doc': conv_doc,
    '.xlsx': conv_xlsx,
    '.xls': conv_xls,
    '.pptx': conv_pptx,
    '.ppt': conv_ppt,
    '.pdf': conv_pdf,
}

if __name__ == '__main__':
    path = sys.argv[1]
    ext = os.path.splitext(path)[1].lower()
    fn = mapping.get(ext)
    if fn:
        try:
            text = fn(path)
            sys.stdout.write(text)
        except Exception as e:
            pass
    else:
        with open(path, 'r', errors='ignore') as f:
            sys.stdout.write(f.read())
